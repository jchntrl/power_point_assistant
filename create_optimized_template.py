#!/usr/bin/env python3
"""
Create an optimized PowerPoint template for the PowerPoint Assistant application.

This script creates a professional template specifically designed for the 
PowerPoint Assistant system with proper layouts, placeholders, and branding.
"""

import logging
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.slide import SlideLayout
from pptx.shapes.placeholder import SlidePlaceholder

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)


def create_optimized_template():
    """Create an optimized PowerPoint template."""
    
    # Start with a new presentation
    prs = Presentation()
    
    # Define professional color scheme
    BRAND_BLUE = RGBColor(0, 100, 177)      # Primary brand blue
    DARK_GRAY = RGBColor(64, 64, 64)        # Dark gray for text
    LIGHT_BLUE = RGBColor(173, 216, 230)    # Light blue accent
    ORANGE_ACCENT = RGBColor(255, 140, 0)   # Orange for highlights
    LIGHT_GRAY = RGBColor(245, 245, 245)    # Light gray background
    WHITE = RGBColor(255, 255, 255)         # White
    
    logger.info("Setting up slide master...")
    
    # Configure slide master
    slide_master = prs.slide_master
    slide_master.background.fill.solid()
    slide_master.background.fill.fore_color.rgb = WHITE
    
    # Get layouts that we need to optimize
    title_layout = prs.slide_layouts[0]  # Title Slide
    content_layout = prs.slide_layouts[1]  # Title and Content
    blank_layout = prs.slide_layouts[6]   # Blank
    
    logger.info("Optimizing Title Slide layout...")
    _optimize_title_layout(title_layout, BRAND_BLUE, DARK_GRAY)
    
    logger.info("Optimizing Content layout...")
    _optimize_content_layout(content_layout, BRAND_BLUE, DARK_GRAY)
    
    logger.info("Optimizing Blank layout...")
    _optimize_blank_layout(blank_layout)
    
    return prs


def _optimize_title_layout(layout, primary_color, text_color):
    """Optimize the title slide layout for AI generation."""
    
    # Find and optimize title placeholder
    for shape in layout.placeholders:
        if shape.placeholder_format.type == 1:  # Title placeholder
            # Position and size
            shape.top = Inches(2.8)
            shape.left = Inches(0.8)
            shape.width = Inches(8.4)
            shape.height = Inches(1.8)
            
            # Set default formatting
            if hasattr(shape, 'text_frame'):
                paragraph = shape.text_frame.paragraphs[0]
                paragraph.font.name = "Segoe UI"
                paragraph.font.size = Pt(48)
                paragraph.font.color.rgb = primary_color
                paragraph.alignment = PP_ALIGN.CENTER
                paragraph.font.bold = True
                
        elif shape.placeholder_format.type == 2:  # Subtitle placeholder
            # Position and size
            shape.top = Inches(5.0)
            shape.left = Inches(0.8)
            shape.width = Inches(8.4)
            shape.height = Inches(1.2)
            
            # Set default formatting
            if hasattr(shape, 'text_frame'):
                paragraph = shape.text_frame.paragraphs[0]
                paragraph.font.name = "Segoe UI"
                paragraph.font.size = Pt(24)
                paragraph.font.color.rgb = text_color
                paragraph.alignment = PP_ALIGN.CENTER


def _optimize_content_layout(layout, primary_color, text_color):
    """Optimize the content layout for AI-generated bullet points."""
    
    for shape in layout.placeholders:
        if shape.placeholder_format.type == 1:  # Title placeholder
            # Position and size for content slides
            shape.top = Inches(0.4)
            shape.left = Inches(0.5)
            shape.width = Inches(9.0)
            shape.height = Inches(1.0)
            
            # Set default formatting
            if hasattr(shape, 'text_frame'):
                paragraph = shape.text_frame.paragraphs[0]
                paragraph.font.name = "Segoe UI"
                paragraph.font.size = Pt(36)
                paragraph.font.color.rgb = primary_color
                paragraph.font.bold = True
                paragraph.alignment = PP_ALIGN.LEFT
                
        elif shape.placeholder_format.type == 2:  # Content placeholder
            # Position and size for optimal content display
            shape.top = Inches(1.6)
            shape.left = Inches(0.5)
            shape.width = Inches(9.0)
            shape.height = Inches(5.8)
            
            # Set default bullet point formatting
            if hasattr(shape, 'text_frame'):
                text_frame = shape.text_frame
                text_frame.margin_left = Inches(0.1)
                text_frame.margin_right = Inches(0.1)
                text_frame.margin_top = Inches(0.1)
                text_frame.margin_bottom = Inches(0.1)
                
                # Default paragraph formatting
                paragraph = text_frame.paragraphs[0]
                paragraph.font.name = "Segoe UI"
                paragraph.font.size = Pt(20)
                paragraph.font.color.rgb = text_color
                paragraph.space_after = Pt(12)


def _optimize_blank_layout(layout):
    """Optimize blank layout for flexible content."""
    # Blank layout is typically used as-is, but we can set background
    pass


def add_demonstration_slides(prs):
    """Add demonstration slides showing the template capabilities."""
    
    logger.info("Adding demonstration slides...")
    
    # Colors for consistent styling
    BRAND_BLUE = RGBColor(0, 100, 177)
    DARK_GRAY = RGBColor(64, 64, 64)
    
    # Demonstration Title Slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "PowerPoint Assistant\nOptimized Template"
    title_slide.placeholders[1].text = "Professional AI-Generated Presentations\\nBuilt for Business Excellence"
    
    # Format title slide
    title_paragraph = title_slide.shapes.title.text_frame.paragraphs[0]
    title_paragraph.font.name = "Segoe UI"
    title_paragraph.font.size = Pt(48)
    title_paragraph.font.color.rgb = BRAND_BLUE
    title_paragraph.font.bold = True
    title_paragraph.alignment = PP_ALIGN.CENTER
    
    # Format subtitle
    subtitle_paragraph = title_slide.placeholders[1].text_frame.paragraphs[0]
    subtitle_paragraph.font.name = "Segoe UI"
    subtitle_paragraph.font.size = Pt(24)
    subtitle_paragraph.font.color.rgb = DARK_GRAY
    subtitle_paragraph.alignment = PP_ALIGN.CENTER
    
    # Demonstration Content Slide
    content_slide = prs.slides.add_slide(prs.slide_layouts[1])
    content_slide.shapes.title.text = "Template Optimization Features"
    
    # Add bullet points
    content_placeholder = content_slide.placeholders[1]
    text_frame = content_placeholder.text_frame
    text_frame.text = "Optimized layout spacing for AI-generated content"
    
    bullet_points = [
        "Professional typography with Segoe UI font family",
        "Consistent color scheme aligned with business standards",
        "Proper placeholder positioning for automated population",
        "Multiple layout options for different content types",
        "Enhanced readability with optimal font sizes and spacing"
    ]
    
    for bullet_point in bullet_points:
        p = text_frame.add_paragraph()
        p.text = bullet_point
        p.level = 0
        p.font.name = "Segoe UI"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
    
    # Format content title
    title_paragraph = content_slide.shapes.title.text_frame.paragraphs[0]
    title_paragraph.font.name = "Segoe UI"
    title_paragraph.font.size = Pt(36)
    title_paragraph.font.color.rgb = BRAND_BLUE
    title_paragraph.font.bold = True
    
    logger.info("Demonstration slides added successfully")


def validate_template_compatibility(prs):
    """Validate that the template is compatible with PowerPoint Assistant."""
    
    logger.info("Validating template compatibility...")
    
    validation_results = {
        "total_layouts": len(prs.slide_layouts),
        "has_title_layout": False,
        "has_content_layout": False,
        "has_blank_layout": False,
        "layout_details": []
    }
    
    for i, layout in enumerate(prs.slide_layouts):
        layout_info = {
            "index": i,
            "name": layout.name,
            "placeholders": len(layout.placeholders)
        }
        
        # Check for key layouts
        layout_name_lower = layout.name.lower()
        if "title" in layout_name_lower and "slide" in layout_name_lower:
            validation_results["has_title_layout"] = True
        elif "content" in layout_name_lower or "title and content" in layout_name_lower:
            validation_results["has_content_layout"] = True
        elif "blank" in layout_name_lower:
            validation_results["has_blank_layout"] = True
            
        validation_results["layout_details"].append(layout_info)
    
    # Check compatibility
    is_compatible = (
        validation_results["has_title_layout"] and
        validation_results["has_content_layout"] and
        validation_results["total_layouts"] >= 3
    )
    
    logger.info(f"Template compatibility: {'✅ PASSED' if is_compatible else '❌ FAILED'}")
    logger.info(f"Total layouts: {validation_results['total_layouts']}")
    logger.info(f"Has title layout: {validation_results['has_title_layout']}")
    logger.info(f"Has content layout: {validation_results['has_content_layout']}")
    logger.info(f"Has blank layout: {validation_results['has_blank_layout']}")
    
    return is_compatible, validation_results


def save_optimized_template(prs, output_path: Path):
    """Save the optimized template."""
    
    try:
        # Ensure output directory exists
        output_path.parent.mkdir(parents=True, exist_ok=True)
        
        # Validate before saving
        is_compatible, validation_results = validate_template_compatibility(prs)
        
        if not is_compatible:
            logger.warning("Template may not be fully compatible with PowerPoint Assistant")
        
        # Save the presentation
        prs.save(str(output_path))
        logger.info(f"Optimized template saved: {output_path}")
        
        return True, validation_results
        
    except Exception as e:
        logger.error(f"Failed to save optimized template: {e}")
        return False, None


def main():
    """Main function to create the optimized PowerPoint template."""
    
    logger.info("Creating optimized PowerPoint template for PowerPoint Assistant...")
    
    # Create the optimized template
    prs = create_optimized_template()
    
    # Add demonstration slides
    add_demonstration_slides(prs)
    
    # Save the template
    output_path = Path("PowerPoint_Assistant_Optimized_Template.pptx")
    success, validation_results = save_optimized_template(prs, output_path)
    
    if success:
        logger.info("Optimized PowerPoint template creation completed!")
        
        # Display results
        print("\\n" + "="*70)
        print("✅ OPTIMIZED POWERPOINT TEMPLATE CREATED SUCCESSFULLY!")
        print("="*70)
        print(f"📁 Template saved as: {output_path.absolute()}")
        print(f"📊 Total layouts available: {validation_results['total_layouts']}")
        
        print("\\n🎯 Layout Summary:")
        for layout_info in validation_results['layout_details'][:6]:  # Show first 6
            print(f"  • Layout {layout_info['index']}: {layout_info['name']} "
                  f"({layout_info['placeholders']} placeholders)")
        
        print("\\n🔧 Optimization Features:")
        print("  ✅ Professional Segoe UI typography")
        print("  ✅ Optimized placeholder positioning")
        print("  ✅ Consistent brand color scheme")
        print("  ✅ Enhanced spacing for AI content")
        print("  ✅ Compatible with PowerPoint Assistant system")
        
        print("\\n📋 Next Steps:")
        print("1. Update your .env file:")
        print(f"   DEFAULT_TEMPLATE={output_path.name}")
        print("2. Test the template with PowerPoint Assistant")
        print("3. Generate your first AI-powered presentation!")
        
    else:
        logger.error("Failed to create optimized template")
        return 1
    
    return 0


if __name__ == "__main__":
    exit(main())
"""
Presentation builder tool for creating PowerPoint files.

This module creates final PowerPoint presentations from slide specifications
using company templates with proper branding and formatting.
"""

import logging
from datetime import datetime
from pathlib import Path
from typing import Optional

from ..config.settings import settings
from ..models.data_models import (
    ContentGenerationResult,
    GeneratedDiagram,
    GeneratedSlide,
    PresentationSpec,
    ProjectDescription,
)
from .template_manager import TemplateManager

logger = logging.getLogger(__name__)


class PresentationBuilder:
    """
    Builds PowerPoint presentations from generated content.
    
    Uses template manager to create professional presentations with
    proper branding and consistent formatting.
    """

    def __init__(self, template_path: Optional[Path] = None) -> None:
        """
        Initialize the presentation builder.

        Args:
            template_path: Optional path to specific template file
        """
        self.template_manager = TemplateManager(template_path)
        self.current_spec = None

    async def build_presentation(
        self,
        project: ProjectDescription,
        generation_result: ContentGenerationResult,
        output_filename: Optional[str] = None,
        template_path: Optional[Path] = None
    ) -> Path:
        """
        Build complete PowerPoint presentation.

        Args:
            project: Project description and details
            generation_result: Generated content from content generation chain
            output_filename: Optional custom output filename
            template_path: Optional template file path

        Returns:
            Path to created presentation file

        Raises:
            ValueError: If presentation building fails
        """
        try:
            logger.info(f"Building presentation for {project.client_name}")

            # Generate output filename if not provided
            if not output_filename:
                output_filename = self._generate_output_filename(project)

            # Ensure output path has correct extension
            if not output_filename.endswith('.pptx'):
                output_filename += '.pptx'

            # Create full output path
            output_path = settings.output_dir / output_filename

            # Create presentation specification
            presentation_spec = PresentationSpec(
                project=project,
                slides=generation_result.slides,
                template_path=template_path or settings.template_path,
                output_path=output_path
            )

            # Store current spec for potential adjustments
            self.current_spec = presentation_spec

            # Validate template
            is_valid, error_msg = self.template_manager.validate_template(presentation_spec.template_path)
            if not is_valid:
                logger.warning(f"Template validation failed: {error_msg}")
                # Continue with default template handling

            # Create the presentation
            created_path = self.template_manager.create_presentation_from_spec(presentation_spec)

            # Add metadata and properties
            self._set_presentation_properties(created_path, project, generation_result)

            logger.info(f"Successfully created presentation: {created_path}")
            return created_path

        except Exception as e:
            logger.error(f"Failed to build presentation: {e}")
            raise ValueError(f"Presentation building failed: {e}") from e

    def _generate_output_filename(self, project: ProjectDescription) -> str:
        """
        Generate appropriate output filename.

        Args:
            project: Project description

        Returns:
            Generated filename string
        """
        # Clean client name for filename
        client_clean = "".join(c for c in project.client_name if c.isalnum() or c in (' ', '-', '_'))
        client_clean = client_clean.strip().replace(' ', '_')

        # Add timestamp
        timestamp = datetime.now().strftime("%Y%m%d_%H%M")

        # Create filename
        filename = f"{client_clean}_Proposal_{timestamp}.pptx"

        return filename

    def _set_presentation_properties(
        self,
        presentation_path: Path,
        project: ProjectDescription,
        generation_result: ContentGenerationResult
    ) -> None:
        """
        Set presentation properties and metadata.

        Args:
            presentation_path: Path to the presentation file
            project: Project description
            generation_result: Generation result with metadata
        """
        try:
            from pptx import Presentation

            # Load the presentation
            prs = Presentation(str(presentation_path))

            # Set core properties
            core_props = prs.core_properties
            core_props.title = f"{project.description[:50]}... - Proposal for {project.client_name}"
            core_props.author = "Keyrus"
            core_props.subject = f"Technology Proposal for {project.client_name}"
            core_props.comments = f"Generated presentation with {len(generation_result.slides)} slides"
            core_props.category = "Technology Proposal"
            core_props.keywords = f"Keyrus, {project.client_name}, Technology, Proposal"
            core_props.created = datetime.now()
            core_props.modified = datetime.now()

            # Save updated properties
            prs.save(str(presentation_path))

        except Exception as e:
            logger.warning(f"Failed to set presentation properties: {e}")

    def customize_slide_content(
        self,
        slide_index: int,
        new_content: list[str],
        new_title: Optional[str] = None,
        new_notes: Optional[str] = None
    ) -> bool:
        """
        Customize content of a specific slide.

        Args:
            slide_index: Index of slide to customize (0-based)
            new_content: New bullet point content
            new_title: Optional new title
            new_notes: Optional new speaker notes

        Returns:
            True if customization successful, False otherwise
        """
        if not self.current_spec or slide_index >= len(self.current_spec.slides):
            logger.error(f"Invalid slide index: {slide_index}")
            return False

        try:
            slide = self.current_spec.slides[slide_index]

            # Update content
            slide.content = new_content

            if new_title:
                slide.title = new_title

            if new_notes:
                slide.notes = new_notes

            logger.info(f"Customized slide {slide_index}: {slide.title}")
            return True

        except Exception as e:
            logger.error(f"Failed to customize slide {slide_index}: {e}")
            return False

    def add_slide(
        self,
        title: str,
        content: list[str],
        layout_type: str = "bullet",
        notes: Optional[str] = None,
        position: Optional[int] = None
    ) -> bool:
        """
        Add a new slide to the presentation.

        Args:
            title: Slide title
            content: Bullet point content
            layout_type: Slide layout type
            notes: Optional speaker notes
            position: Optional position to insert slide (None = append)

        Returns:
            True if slide added successfully, False otherwise
        """
        if not self.current_spec:
            logger.error("No current presentation specification")
            return False

        try:
            new_slide = GeneratedSlide(
                title=title,
                content=content,
                layout_type=layout_type,
                notes=notes or ""
            )

            if position is not None and 0 <= position <= len(self.current_spec.slides):
                self.current_spec.slides.insert(position, new_slide)
            else:
                self.current_spec.slides.append(new_slide)

            logger.info(f"Added slide: {title}")
            return True

        except Exception as e:
            logger.error(f"Failed to add slide: {e}")
            return False

    def remove_slide(self, slide_index: int) -> bool:
        """
        Remove a slide from the presentation.

        Args:
            slide_index: Index of slide to remove (0-based)

        Returns:
            True if slide removed successfully, False otherwise
        """
        if not self.current_spec or slide_index >= len(self.current_spec.slides):
            logger.error(f"Invalid slide index: {slide_index}")
            return False

        try:
            removed_slide = self.current_spec.slides.pop(slide_index)
            logger.info(f"Removed slide: {removed_slide.title}")
            return True

        except Exception as e:
            logger.error(f"Failed to remove slide {slide_index}: {e}")
            return False

    def reorder_slides(self, new_order: list[int]) -> bool:
        """
        Reorder slides in the presentation.

        Args:
            new_order: List of slide indices in new order

        Returns:
            True if reordering successful, False otherwise
        """
        if not self.current_spec:
            logger.error("No current presentation specification")
            return False

        if len(new_order) != len(self.current_spec.slides):
            logger.error("New order list length doesn't match slide count")
            return False

        if set(new_order) != set(range(len(self.current_spec.slides))):
            logger.error("Invalid slide indices in new order")
            return False

        try:
            original_slides = self.current_spec.slides.copy()
            self.current_spec.slides = [original_slides[i] for i in new_order]

            logger.info("Successfully reordered slides")
            return True

        except Exception as e:
            logger.error(f"Failed to reorder slides: {e}")
            return False

    def rebuild_presentation(self) -> Optional[Path]:
        """
        Rebuild the presentation with current specifications.

        Returns:
            Path to rebuilt presentation or None if failed
        """
        if not self.current_spec:
            logger.error("No current presentation specification")
            return None

        try:
            # Create new presentation with updated specs
            rebuilt_path = self.template_manager.create_presentation_from_spec(self.current_spec)

            logger.info(f"Successfully rebuilt presentation: {rebuilt_path}")
            return rebuilt_path

        except Exception as e:
            logger.error(f"Failed to rebuild presentation: {e}")
            return None

    async def add_diagram_as_dedicated_slide(
        self,
        presentation: any,
        diagram: GeneratedDiagram
    ) -> bool:
        """
        Add diagram as a dedicated slide in the presentation.

        Args:
            presentation: PowerPoint presentation object
            diagram: Generated diagram with image path and specifications

        Returns:
            True if diagram slide created successfully, False otherwise
        """
        try:
            from pptx.util import Inches

            # Verify diagram image exists
            if not diagram.image_path.exists():
                logger.error(f"Diagram image not found: {diagram.image_path}")
                return False

            # Get optimal layout for diagram
            layout_type = self.template_manager.get_optimal_layout_for_diagram(
                diagram.spec.diagram_type,
                has_content=False  # Dedicated slide has no text content
            )

            # Get layout index
            layout_index = self.template_manager.get_layout_index(layout_type)
            layout = presentation.slide_layouts[layout_index]

            # Create new slide
            slide = presentation.slides.add_slide(layout)

            # Set slide title
            if slide.shapes.title:
                slide.shapes.title.text = diagram.spec.title

            # Calculate optimal diagram positioning for full slide
            # Use nearly full slide dimensions with some padding

            # Position diagram to take up most of the slide
            left = Inches(0.5)  # Small left margin
            top = Inches(1.5)   # Leave space for title
            width = Inches(9.0)  # Almost full width
            height = Inches(5.5)  # Remaining height after title

            # Add diagram image to slide
            pic = slide.shapes.add_picture(
                str(diagram.image_path),
                left, top, width, height
            )

            # Verify image was added successfully
            if not pic:
                logger.warning(f"Failed to add diagram '{diagram.spec.title}' to dedicated slide")
                return False

            # Add optional description as speaker notes
            if hasattr(diagram.spec, 'description') and diagram.spec.description:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = diagram.spec.description

            logger.info(
                f"Successfully created dedicated slide for diagram '{diagram.spec.title}' "
                f"with layout '{layout_type}' at position ({left.inches}, {top.inches}) "
                f"with size ({width.inches} x {height.inches})"
            )
            return True

        except Exception as e:
            logger.error(f"Error creating dedicated slide for diagram '{diagram.spec.title}': {e}")
            return False

    async def insert_diagrams_into_presentation(
        self,
        presentation_path: Path,
        diagrams: list[GeneratedDiagram]
    ) -> dict[str, any]:
        """
        Insert multiple diagrams as dedicated pages in the presentation.

        Args:
            presentation_path: Path to the presentation file
            diagrams: List of generated diagrams to insert

        Returns:
            Dictionary with insertion results
        """
        results = {
            "total_diagrams": len(diagrams),
            "successful_insertions": 0,
            "failed_insertions": 0,
            "insertion_details": []
        }

        if not diagrams:
            logger.info("No diagrams to insert")
            return results

        try:
            from pptx import Presentation

            # Load presentation
            prs = Presentation(str(presentation_path))

            for diagram in diagrams:
                try:
                    # Create a dedicated slide for this diagram
                    success = await self.add_diagram_as_dedicated_slide(
                        prs, diagram
                    )

                    if success:
                        results["successful_insertions"] += 1
                        results["insertion_details"].append({
                            "diagram_title": diagram.spec.title,
                            "slide_index": len(prs.slides) - 1,  # Last added slide
                            "status": "success",
                            "image_path": str(diagram.image_path),
                            "file_size_kb": diagram.file_size_kb
                        })
                    else:
                        results["failed_insertions"] += 1
                        results["insertion_details"].append({
                            "diagram_title": diagram.spec.title,
                            "slide_index": -1,
                            "status": "failed",
                            "error": "Failed to create dedicated diagram slide"
                        })

                except Exception as e:
                    results["failed_insertions"] += 1
                    results["insertion_details"].append({
                        "diagram_title": diagram.spec.title,
                        "slide_index": -1,
                        "status": "failed",
                        "error": str(e)
                    })
                    logger.error(f"Failed to insert diagram '{diagram.spec.title}': {e}")

            # Save the presentation with new diagram slides
            prs.save(str(presentation_path))

            logger.info(
                f"Diagram insertion complete: {results['successful_insertions']} successful, "
                f"{results['failed_insertions']} failed"
            )

        except Exception as e:
            logger.error(f"Error during diagram insertion: {e}")
            results["error"] = str(e)

        return results

    def _optimize_diagram_positioning(
        self,
        slide_layout_name: str,
        diagram_type: str
    ) -> dict[str, float]:
        """
        Optimize diagram positioning based on slide layout and diagram type.

        Args:
            slide_layout_name: Name of the slide layout
            diagram_type: Type of diagram being positioned

        Returns:
            Dictionary with optimized position values
        """
        # Define layout-specific positioning
        layout_positions = {
            "Title Slide": {"left": 1.0, "top": 3.0, "width": 8.0, "height": 4.0},
            "Title and Content": {"left": 0.5, "top": 1.8, "width": 9.0, "height": 5.0},
            "Content with Caption": {"left": 0.5, "top": 1.5, "width": 6.0, "height": 5.5},
            "Two Content": {"left": 5.0, "top": 1.5, "width": 4.5, "height": 5.5},
            "Blank": {"left": 0.5, "top": 1.0, "width": 9.0, "height": 6.5}
        }

        # Get base position for layout
        base_position = layout_positions.get(
            slide_layout_name,
            layout_positions["Title and Content"]
        )

        # Adjust for diagram type
        adjustments = {
            "microservices": {"height": 5.5},
            "data_pipeline": {"width": 9.5, "height": 4.5},
            "cloud_architecture": {"width": 9.0, "height": 6.0},
            "database_schema": {"width": 7.0, "height": 5.0}
        }

        if diagram_type in adjustments:
            position = dict(base_position)
            position.update(adjustments[diagram_type])
            return position

        return base_position

    def get_presentation_summary(self) -> dict:
        """
        Get summary information about the current presentation.

        Returns:
            Dictionary with presentation summary
        """
        if not self.current_spec:
            return {"error": "No current presentation specification"}

        slide_summaries = []
        for i, slide in enumerate(self.current_spec.slides):
            slide_summaries.append({
                "index": i,
                "title": slide.title,
                "content_count": len(slide.content),
                "layout_type": slide.layout_type,
                "has_notes": bool(slide.notes)
            })

        return {
            "client_name": self.current_spec.project.client_name,
            "total_slides": len(self.current_spec.slides),
            "template_path": str(self.current_spec.template_path),
            "output_path": str(self.current_spec.output_path),
            "slides": slide_summaries
        }

    def validate_presentation_structure(self) -> dict:
        """
        Validate the current presentation structure.

        Returns:
            Dictionary with validation results
        """
        if not self.current_spec:
            return {"valid": False, "errors": ["No presentation specification"]}

        errors = []
        warnings = []

        # Check slide count
        slide_count = len(self.current_spec.slides)
        if slide_count < settings.min_slides:
            errors.append(f"Too few slides: {slide_count} (minimum: {settings.min_slides})")
        elif slide_count > settings.max_slides:
            warnings.append(f"Many slides: {slide_count} (maximum recommended: {settings.max_slides})")

        # Check individual slides
        for i, slide in enumerate(self.current_spec.slides):
            # Check title
            if not slide.title or len(slide.title.strip()) < 3:
                errors.append(f"Slide {i+1}: Title too short or missing")

            # Check content
            if not slide.content or len(slide.content) == 0:
                errors.append(f"Slide {i+1}: No content")
            elif len(slide.content) > 7:
                warnings.append(f"Slide {i+1}: Too many bullet points ({len(slide.content)})")

            # Check for very long bullet points
            for j, bullet in enumerate(slide.content):
                if len(bullet) > 150:
                    warnings.append(f"Slide {i+1}, bullet {j+1}: Very long bullet point")

        # Check for duplicate titles
        titles = [slide.title for slide in self.current_spec.slides]
        if len(titles) != len(set(titles)):
            warnings.append("Duplicate slide titles found")

        return {
            "valid": len(errors) == 0,
            "errors": errors,
            "warnings": warnings,
            "slide_count": slide_count,
            "total_issues": len(errors) + len(warnings)
        }


# Convenience function for simple presentation building
async def build_presentation_simple(
    project: ProjectDescription,
    generation_result: ContentGenerationResult,
    output_filename: Optional[str] = None
) -> Path:
    """
    Simple function to build presentation without creating a builder instance.

    Args:
        project: Project description
        generation_result: Generated content
        output_filename: Optional output filename

    Returns:
        Path to created presentation
    """
    builder = PresentationBuilder()
    return await builder.build_presentation(project, generation_result, output_filename)

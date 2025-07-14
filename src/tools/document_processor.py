"""
Document processing tools for PowerPoint and PDF files.

This module provides utilities to extract text content from PowerPoint presentations
and PDF documents, handling various file formats and edge cases safely.
"""

import logging
from io import BytesIO
from pathlib import Path
from typing import List, Optional, Union

import pypdf
from pptx import Presentation
from pptx.exc import PackageNotFoundError

from ..models.data_models import ExtractedContent

logger = logging.getLogger(__name__)


class DocumentProcessor:
    """
    Processor for extracting text content from PowerPoint and PDF files.
    
    Handles both file paths and BytesIO objects (for Streamlit uploads).
    """

    def __init__(self) -> None:
        """Initialize the document processor."""
        self.supported_extensions = {".pptx", ".pdf"}

    async def process_document(
        self,
        file_source: Union[Path, BytesIO],
        filename: str,
        file_type: str
    ) -> List[ExtractedContent]:
        """
        Process a document and extract its content.

        Args:
            file_source: Path to file or BytesIO object with file content
            filename: Original filename for metadata
            file_type: File type ("pptx" or "pdf")

        Returns:
            List of extracted content from the document

        Raises:
            ValueError: If file type is not supported
            Exception: If document processing fails
        """
        file_type = file_type.lower()
        
        if file_type not in ["pptx", "pdf"]:
            raise ValueError(f"Unsupported file type: {file_type}")

        try:
            if file_type == "pptx":
                return await self._extract_from_powerpoint(file_source, filename)
            elif file_type == "pdf":
                return await self._extract_from_pdf(file_source, filename)
        except Exception as e:
            logger.error(f"Failed to process {filename}: {e}")
            raise

    async def _extract_from_powerpoint(
        self,
        file_source: Union[Path, BytesIO],
        filename: str
    ) -> List[ExtractedContent]:
        """
        Extract text content from PowerPoint file.

        Args:
            file_source: Path to file or BytesIO object
            filename: Original filename for metadata

        Returns:
            List of ExtractedContent objects
        """
        try:
            # PATTERN: Use python-pptx with careful error handling
            presentation = Presentation(file_source)
            extracted_content = []

            for slide_num, slide in enumerate(presentation.slides, 1):
                # GOTCHA: Handle different slide layouts safely
                title = ""
                content = ""

                # CRITICAL: Use try-except for placeholder access
                try:
                    if slide.shapes.title:
                        title = slide.shapes.title.text.strip()
                except AttributeError:
                    title = f"Slide {slide_num}"

                # Extract text from all text shapes
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content.append(shape.text.strip())

                content = "\n".join(text_content)

                # Get layout name safely
                layout_type = "unknown"
                try:
                    layout_type = slide.slide_layout.name
                except AttributeError:
                    layout_type = f"layout_{slide.slide_layout_id}"

                # Only add slides with meaningful content
                if content.strip() or title.strip():
                    extracted_content.append(ExtractedContent(
                        slide_number=slide_num,
                        title=title or f"Slide {slide_num}",
                        content=content,
                        layout_type=layout_type,
                        source_file=filename,
                        file_type="pptx"
                    ))

            logger.info(f"Extracted {len(extracted_content)} slides from {filename}")
            return extracted_content

        except PackageNotFoundError:
            logger.error(f"Invalid PowerPoint file: {filename}")
            raise ValueError(f"File {filename} is not a valid PowerPoint presentation")
        except Exception as e:
            logger.error(f"Error extracting from PowerPoint {filename}: {e}")
            raise

    async def _extract_from_pdf(
        self,
        file_source: Union[Path, BytesIO],
        filename: str
    ) -> List[ExtractedContent]:
        """
        Extract text content from PDF file.

        Args:
            file_source: Path to file or BytesIO object
            filename: Original filename for metadata

        Returns:
            List of ExtractedContent objects
        """
        try:
            # Handle BytesIO vs Path
            if isinstance(file_source, BytesIO):
                pdf_reader = pypdf.PdfReader(file_source)
            else:
                pdf_reader = pypdf.PdfReader(str(file_source))

            extracted_content = []

            for page_num, page in enumerate(pdf_reader.pages, 1):
                try:
                    # Extract text from page
                    text = page.extract_text()
                    
                    if not text.strip():
                        continue

                    # Try to extract a title from the first line or first few words
                    lines = text.strip().split('\n')
                    title = ""
                    content = text.strip()

                    if lines:
                        # Use first non-empty line as title if it's short enough
                        first_line = lines[0].strip()
                        if len(first_line) <= 100 and len(lines) > 1:
                            title = first_line
                            content = '\n'.join(lines[1:]).strip()
                        else:
                            # Use page number as title
                            title = f"Page {page_num}"

                    extracted_content.append(ExtractedContent(
                        slide_number=page_num,
                        title=title or f"Page {page_num}",
                        content=content,
                        layout_type="pdf_page",
                        source_file=filename,
                        file_type="pdf"
                    ))

                except Exception as e:
                    logger.warning(f"Error extracting page {page_num} from {filename}: {e}")
                    continue

            logger.info(f"Extracted {len(extracted_content)} pages from {filename}")
            return extracted_content

        except Exception as e:
            logger.error(f"Error extracting from PDF {filename}: {e}")
            raise ValueError(f"Failed to process PDF file {filename}: {e}")

    def validate_file_type(self, filename: str) -> bool:
        """
        Validate if the file type is supported.

        Args:
            filename: Name of the file to validate

        Returns:
            True if file type is supported, False otherwise
        """
        file_extension = Path(filename).suffix.lower()
        return file_extension in self.supported_extensions

    def get_file_type(self, filename: str) -> Optional[str]:
        """
        Get the file type from filename.

        Args:
            filename: Name of the file

        Returns:
            File type ("pptx" or "pdf") or None if not supported
        """
        file_extension = Path(filename).suffix.lower()
        
        if file_extension == ".pptx":
            return "pptx"
        elif file_extension == ".pdf":
            return "pdf"
        else:
            return None

    async def process_multiple_documents(
        self,
        documents: List[tuple]
    ) -> List[ExtractedContent]:
        """
        Process multiple documents concurrently.

        Args:
            documents: List of tuples (file_source, filename, file_type)

        Returns:
            Combined list of all extracted content
        """
        all_content = []
        
        for file_source, filename, file_type in documents:
            try:
                content = await self.process_document(file_source, filename, file_type)
                all_content.extend(content)
            except Exception as e:
                logger.error(f"Failed to process {filename}: {e}")
                # Continue processing other documents
                continue

        logger.info(f"Processed {len(documents)} documents, extracted {len(all_content)} content items")
        return all_content

    def estimate_processing_time(self, file_sizes: List[int]) -> float:
        """
        Estimate processing time for files based on their sizes.

        Args:
            file_sizes: List of file sizes in bytes

        Returns:
            Estimated processing time in seconds
        """
        # Rough estimates: 1MB PowerPoint takes ~2 seconds, 1MB PDF takes ~1 second
        total_mb = sum(file_sizes) / (1024 * 1024)
        estimated_seconds = total_mb * 1.5  # Average estimate
        
        return max(5.0, estimated_seconds)  # Minimum 5 seconds


# Convenience functions for backward compatibility
async def extract_text_from_pptx(file_path: Path) -> List[ExtractedContent]:
    """
    Legacy function for extracting text from PowerPoint files.

    Args:
        file_path: Path to PowerPoint file

    Returns:
        List of extracted content
    """
    processor = DocumentProcessor()
    return await processor._extract_from_powerpoint(file_path, file_path.name)


async def extract_text_from_pdf(file_path: Path) -> List[ExtractedContent]:
    """
    Legacy function for extracting text from PDF files.

    Args:
        file_path: Path to PDF file

    Returns:
        List of extracted content
    """
    processor = DocumentProcessor()
    return await processor._extract_from_pdf(file_path, file_path.name)
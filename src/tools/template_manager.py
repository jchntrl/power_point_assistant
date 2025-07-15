"""
PowerPoint template management and slide creation utilities.

This module handles loading Keyrus PowerPoint templates and creating
slides with proper placeholder management and branding preservation.
"""

import logging
from pathlib import Path
from typing import Optional

from pptx import Presentation
from pptx.exc import PackageNotFoundError
from pptx.slide import Slide
from pptx.util import Inches

from ..config.settings import settings
from ..models.data_models import GeneratedSlide, PresentationSpec

logger = logging.getLogger(__name__)


class TemplateManager:
    """
    Manages PowerPoint templates and slide creation.
    
    Handles Keyrus template loading, layout management, and proper
    placeholder positioning to avoid common python-pptx gotchas.
    """

    def __init__(self, template_path: Optional[Path] = None) -> None:
        """
        Initialize the template manager.

        Args:
            template_path: Optional path to specific template file
        """
        self.template_path = template_path or settings.template_path
        self.presentation = None
        self.slide_layouts = None
        self._layout_mapping = {}

    def load_template(self, template_path: Optional[Path] = None) -> Presentation:
        """
        Load PowerPoint template.

        Args:
            template_path: Optional path to template file

        Returns:
            Loaded Presentation object

        Raises:
            FileNotFoundError: If template file not found
            ValueError: If template file is invalid
        """
        template_path = template_path or self.template_path

        if not template_path.exists():
            raise FileNotFoundError(f"Template file not found: {template_path}")

        try:
            # PATTERN: Load company template safely
            self.presentation = Presentation(str(template_path))
            self.slide_layouts = self.presentation.slide_layouts

            # Map layouts for easier access
            self._map_slide_layouts()

            logger.info(f"Loaded template: {template_path.name} with {len(self.slide_layouts)} layouts")
            return self.presentation

        except PackageNotFoundError:
            raise ValueError(f"Invalid PowerPoint template: {template_path}") from None
        except Exception as e:
            logger.error(f"Error loading template {template_path}: {e}")
            raise

    def _map_slide_layouts(self) -> None:
        """Map slide layouts by name and common patterns."""
        self._layout_mapping = {}

        for i, layout in enumerate(self.slide_layouts):
            layout_name = layout.name.lower()
            self._layout_mapping[layout_name] = i

            # Common layout patterns
            if "title" in layout_name and "slide" in layout_name:
                self._layout_mapping["title"] = i
            elif "title" in layout_name and "content" in layout_name:
                self._layout_mapping["title_content"] = i
            elif "content" in layout_name or "bullet" in layout_name:
                self._layout_mapping["bullet"] = i
            elif "blank" in layout_name:
                self._layout_mapping["blank"] = i
            elif "section" in layout_name:
                self._layout_mapping["section"] = i
            elif "two content" in layout_name or "comparison" in layout_name:
                self._layout_mapping["two_content"] = i
            elif "picture" in layout_name or "caption" in layout_name:
                self._layout_mapping["picture_caption"] = i

            # Diagram-specific layout patterns
            if "diagram" in layout_name:
                self._layout_mapping["diagram"] = i
            elif "split" in layout_name:
                self._layout_mapping["split"] = i

        logger.debug(f"Mapped {len(self._layout_mapping)} layout patterns")

    def get_layout_index(self, layout_type: str) -> int:
        """
        Get layout index by type.

        Args:
            layout_type: Type of layout (title, bullet, blank, etc.)

        Returns:
            Layout index

        Raises:
            ValueError: If layout type not found
        """
        layout_type = layout_type.lower()

        # Direct mapping
        if layout_type in self._layout_mapping:
            return self._layout_mapping[layout_type]

        # Fallback patterns
        if layout_type == "title":
            return self._layout_mapping.get("title", 0)
        elif layout_type in ["bullet", "content"]:
            return self._layout_mapping.get("bullet", 1)
        elif layout_type == "blank":
            return self._layout_mapping.get("blank", 6)
        elif layout_type == "diagram":
            # Prefer blank layout for diagrams, then two_content, then default
            return self._layout_mapping.get("diagram",
                   self._layout_mapping.get("blank",
                   self._layout_mapping.get("two_content", 6)))
        elif layout_type == "split":
            # For split content with diagram
            return self._layout_mapping.get("split",
                   self._layout_mapping.get("two_content", 1))

        # Default fallback
        logger.warning(f"Layout type '{layout_type}' not found, using default (1)")
        return 1

    def create_slide_from_spec(self, slide_spec: GeneratedSlide) -> Slide:
        """
        Create a slide from specification.

        Args:
            slide_spec: Slide specification with content

        Returns:
            Created slide object

        Raises:
            ValueError: If template not loaded or slide creation fails
        """
        if not self.presentation:
            raise ValueError("Template not loaded. Call load_template() first.")

        try:
            # Get appropriate layout
            layout_index = self.get_layout_index(slide_spec.layout_type)
            layout = self.slide_layouts[layout_index]

            # Add slide to presentation
            slide = self.presentation.slides.add_slide(layout)

            # Populate slide content
            self._populate_slide_content(slide, slide_spec)

            logger.debug(f"Created slide: {slide_spec.title}")
            return slide

        except Exception as e:
            logger.error(f"Error creating slide '{slide_spec.title}': {e}")
            raise

    def _populate_slide_content(self, slide: Slide, slide_spec: GeneratedSlide) -> None:
        """
        Populate slide with content from specification.

        Args:
            slide: Slide object to populate
            slide_spec: Slide specification with content

        CRITICAL: Handle placeholder positioning correctly per PRP gotchas
        """
        try:
            # CRITICAL: Set all positioning properties together
            # Set title if available
            if slide.shapes.title and slide_spec.title:
                slide.shapes.title.text = slide_spec.title

            # Handle content based on layout type
            if slide_spec.layout_type.lower() == "title":
                self._populate_title_slide(slide, slide_spec)
            else:
                self._populate_content_slide(slide, slide_spec)

            # Add speaker notes if provided
            if slide_spec.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_spec.notes

        except Exception as e:
            logger.warning(f"Error populating slide content: {e}")
            # Continue gracefully - slide will have title at least

    def _populate_title_slide(self, slide: Slide, slide_spec: GeneratedSlide) -> None:
        """
        Populate title slide with content.

        Args:
            slide: Slide object
            slide_spec: Slide specification
        """
        try:
            # For title slides, try to use subtitle placeholder
            if len(slide.placeholders) > 1:
                subtitle_placeholder = slide.placeholders[1]
                if hasattr(subtitle_placeholder, 'text_frame'):
                    # Join content as subtitle text
                    subtitle_text = "\n".join(slide_spec.content)
                    subtitle_placeholder.text = subtitle_text

        except (KeyError, AttributeError, IndexError) as e:
            logger.warning(f"Could not set subtitle on title slide: {e}")

    def _populate_content_slide(self, slide: Slide, slide_spec: GeneratedSlide) -> None:
        """
        Populate content slide with bullet points.

        Args:
            slide: Slide object
            slide_spec: Slide specification
        """
        try:
            # Find content placeholder
            content_placeholder = None

            # Try different placeholder indices
            for placeholder_idx in [1, 2]:
                try:
                    if placeholder_idx < len(slide.placeholders):
                        placeholder = slide.placeholders[placeholder_idx]
                        if hasattr(placeholder, 'text_frame'):
                            content_placeholder = placeholder
                            break
                except (KeyError, AttributeError):
                    continue

            if not content_placeholder:
                logger.warning("No content placeholder found on slide")
                return

            # CRITICAL: Handle text frame properly
            text_frame = content_placeholder.text_frame
            text_frame.clear()  # Clear default text

            # Add bullet points
            for i, bullet_point in enumerate(slide_spec.content):
                if i == 0:
                    # Use existing paragraph
                    p = text_frame.paragraphs[0]
                else:
                    # Add new paragraph
                    p = text_frame.add_paragraph()

                p.text = bullet_point
                p.level = 0  # Main bullet level

        except Exception as e:
            logger.warning(f"Error setting content placeholder: {e}")

    def get_optimal_layout_for_diagram(self, diagram_type: str, has_content: bool = True) -> str:
        """
        Get optimal slide layout for a diagram type.

        Args:
            diagram_type: Type of diagram (microservices, data_pipeline, etc.)
            has_content: Whether slide will have text content alongside diagram

        Returns:
            Recommended layout type
        """
        # For dedicated diagram slides (has_content=False), prioritize layouts that give maximum space
        if not has_content:
            # Priority order for dedicated diagram slides
            preferred_layouts = ["blank", "diagram", "title_content", "bullet"]

            for layout in preferred_layouts:
                if layout in self._layout_mapping:
                    logger.debug(f"Selected layout '{layout}' for dedicated diagram '{diagram_type}'")
                    return layout
        else:
            # Diagram-specific layout preferences when there's content
            layout_preferences = {
                "microservices": "split",
                "data_pipeline": "two_content",
                "cloud_architecture": "split",
                "database_schema": "picture_caption"
            }

            preferred_layout = layout_preferences.get(diagram_type, "two_content")

            # Fallback to available layouts if preferred not available
            if preferred_layout not in self._layout_mapping:
                preferred_layout = "two_content" if "two_content" in self._layout_mapping else "bullet"

            logger.debug(f"Selected layout '{preferred_layout}' for diagram type '{diagram_type}' with content")
            return preferred_layout

        # Final fallback
        logger.warning(f"No suitable layout found for diagram type '{diagram_type}', using default")
        return "bullet"

    def adjust_slide_for_diagram(self, slide: Slide, layout_type: str) -> dict[str, any]:
        """
        Adjust slide layout to accommodate diagram insertion.

        Args:
            slide: Slide object to adjust
            layout_type: Type of layout being used

        Returns:
            Dictionary with adjustment information
        """
        adjustment_info = {
            "layout_type": layout_type,
            "adjustments_made": [],
            "content_area": None,
            "diagram_area": None
        }

        try:
            # Get slide dimensions (standard PowerPoint slide)

            if layout_type in ["two_content", "split"]:
                # Split layout - content on left, diagram on right
                adjustment_info["content_area"] = {
                    "left": Inches(0.5),
                    "top": Inches(1.5),
                    "width": Inches(4.5),
                    "height": Inches(5.5)
                }
                adjustment_info["diagram_area"] = {
                    "left": Inches(5.0),
                    "top": Inches(1.5),
                    "width": Inches(4.5),
                    "height": Inches(5.5)
                }
                adjustment_info["adjustments_made"].append("Split layout configured")

            elif layout_type in ["blank", "diagram"]:
                # Full slide for diagram with minimal text
                adjustment_info["content_area"] = {
                    "left": Inches(0.5),
                    "top": Inches(1.0),
                    "width": Inches(9.0),
                    "height": Inches(1.0)
                }
                adjustment_info["diagram_area"] = {
                    "left": Inches(0.5),
                    "top": Inches(2.2),
                    "width": Inches(9.0),
                    "height": Inches(5.0)
                }
                adjustment_info["adjustments_made"].append("Full diagram layout configured")

            elif layout_type == "picture_caption":
                # Diagram with caption below
                adjustment_info["diagram_area"] = {
                    "left": Inches(1.0),
                    "top": Inches(1.5),
                    "width": Inches(8.0),
                    "height": Inches(4.5)
                }
                adjustment_info["content_area"] = {
                    "left": Inches(1.0),
                    "top": Inches(6.0),
                    "width": Inches(8.0),
                    "height": Inches(1.0)
                }
                adjustment_info["adjustments_made"].append("Picture caption layout configured")

            else:
                # Default content layout with space for diagram
                adjustment_info["content_area"] = {
                    "left": Inches(0.5),
                    "top": Inches(1.5),
                    "width": Inches(4.0),
                    "height": Inches(5.5)
                }
                adjustment_info["diagram_area"] = {
                    "left": Inches(5.0),
                    "top": Inches(1.5),
                    "width": Inches(4.5),
                    "height": Inches(5.5)
                }
                adjustment_info["adjustments_made"].append("Default content layout with diagram space")

            logger.debug(f"Adjusted slide for diagram: {adjustment_info['adjustments_made']}")

        except Exception as e:
            logger.error(f"Error adjusting slide for diagram: {e}")
            adjustment_info["error"] = str(e)

        return adjustment_info

    def create_presentation_from_spec(self, spec: PresentationSpec) -> Path:
        """
        Create complete presentation from specification.

        Args:
            spec: Complete presentation specification

        Returns:
            Path to created presentation file

        Raises:
            ValueError: If template loading or slide creation fails
        """
        try:
            # Load template
            self.load_template(spec.template_path)

            # Remove existing slides (except keep master)
            slide_count = len(self.presentation.slides)
            for i in range(slide_count - 1, -1, -1):
                slide_id = self.presentation.slides._sldIdLst[i]
                self.presentation.part.drop_rel(slide_id.rId)
                del self.presentation.slides._sldIdLst[i]

            # Create slides from specification
            for slide_spec in spec.slides:
                self.create_slide_from_spec(slide_spec)

            # Ensure output directory exists
            spec.output_path.parent.mkdir(parents=True, exist_ok=True)

            # Save presentation
            self.presentation.save(str(spec.output_path))

            logger.info(f"Created presentation: {spec.output_path}")
            return spec.output_path

        except Exception as e:
            logger.error(f"Error creating presentation: {e}")
            raise

    def get_template_info(self) -> dict:
        """
        Get information about the loaded template.

        Returns:
            Dictionary with template information
        """
        if not self.presentation:
            return {"error": "No template loaded"}

        return {
            "slide_layouts": len(self.slide_layouts),
            "layout_names": [layout.name for layout in self.slide_layouts],
            "layout_mapping": self._layout_mapping,
            "template_path": str(self.template_path)
        }

    def validate_template(self, template_path: Path) -> tuple[bool, str]:
        """
        Validate a PowerPoint template file.

        Args:
            template_path: Path to template file

        Returns:
            Tuple of (is_valid, error_message)
        """
        try:
            if not template_path.exists():
                return False, f"Template file not found: {template_path}"

            # Try to load the template
            test_presentation = Presentation(str(template_path))

            # Check basic requirements
            if len(test_presentation.slide_layouts) == 0:
                return False, "Template has no slide layouts"

            # Check for common layouts
            layout_names = [layout.name.lower() for layout in test_presentation.slide_layouts]
            has_title_layout = any("title" in name for name in layout_names)
            has_content_layout = any(("content" in name or "bullet" in name) for name in layout_names)

            if not (has_title_layout and has_content_layout):
                return False, "Template missing required layouts (title, content)"

            return True, "Template is valid"

        except Exception as e:
            return False, f"Template validation failed: {e}"

    def list_available_templates(self) -> list[Path]:
        """
        List available PowerPoint templates.

        Returns:
            List of template file paths
        """
        templates = []

        # Check template directory
        if settings.template_dir.exists():
            templates.extend(settings.template_dir.glob("*.pptx"))

        # Check project root for Keyrus templates
        project_root = Path(".")
        templates.extend(project_root.glob("*Template*.pptx"))
        templates.extend(project_root.glob("*template*.pptx"))

        # Remove duplicates and sort
        unique_templates = list(set(templates))
        unique_templates.sort()

        return unique_templates

    def create_diagram_slide_spec(self, diagram) -> 'GeneratedSlide':
        """
        Create a slide specification for a diagram.

        Args:
            diagram: Generated diagram to create slide spec for

        Returns:
            GeneratedSlide specification for the diagram
        """
        from ..models.data_models import GeneratedSlide

        # Get optimal layout for this diagram type
        layout_type = self.get_optimal_layout_for_diagram(
            diagram.spec.diagram_type,
            has_content=False
        )

        # Create slide specification
        slide_spec = GeneratedSlide(
            title=diagram.spec.title,
            content=[],  # No text content for dedicated diagram slides
            layout_type=layout_type,
            notes=f"Architecture diagram showing {diagram.spec.diagram_type} components and their relationships."
        )

        logger.debug(f"Created diagram slide spec: {diagram.spec.title} with layout {layout_type}")
        return slide_spec

    @staticmethod
    def get_default_slide_specs(project_title: str, client_name: str) -> list[GeneratedSlide]:
        """
        Generate default slide specifications for a presentation.

        Args:
            project_title: Title of the project
            client_name: Name of the client

        Returns:
            List of default slide specifications
        """
        return [
            GeneratedSlide(
                title=f"{project_title}",
                content=[f"Proposal for {client_name}", "Prepared by Keyrus"],
                layout_type="title",
                notes="Opening slide with project title and client information"
            ),
            GeneratedSlide(
                title="Executive Summary",
                content=[
                    "Project overview and objectives",
                    "Key benefits and value proposition",
                    "High-level approach and methodology",
                    "Expected outcomes and deliverables"
                ],
                layout_type="bullet",
                notes="High-level overview of the project proposal"
            ),
            GeneratedSlide(
                title="Technical Approach",
                content=[
                    "Solution architecture overview",
                    "Technology stack and tools",
                    "Implementation methodology",
                    "Integration considerations"
                ],
                layout_type="bullet",
                notes="Detailed technical approach and methodology"
            ),
            GeneratedSlide(
                title="Project Timeline",
                content=[
                    "Phase 1: Planning and Design",
                    "Phase 2: Development and Implementation",
                    "Phase 3: Testing and Deployment",
                    "Phase 4: Go-live and Support"
                ],
                layout_type="bullet",
                notes="High-level project phases and timeline"
            ),
            GeneratedSlide(
                title="Next Steps",
                content=[
                    "Project kick-off and team formation",
                    "Detailed requirements gathering",
                    "Technical architecture finalization",
                    "Contract and timeline confirmation"
                ],
                layout_type="bullet",
                notes="Immediate next steps for project initiation"
            )
        ]
